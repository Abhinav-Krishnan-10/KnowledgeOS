import os
import re
import logging
from typing import List, Dict, Any
import fitz  # PyMuPDF
import docx
import pptx

class RecursiveCharacterTextSplitter:
    """A lightweight, zero-dependency recursive character text splitter."""
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200, separators: List[str] = None):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.separators = separators or ["\n\n", "\n", " ", ""]

    def split_text(self, text: str) -> List[str]:
        if not text:
            return []
        
        chunks = []
        
        def recurse(text_to_split: str, separators_list: List[str]):
            if len(text_to_split) <= self.chunk_size:
                chunks.append(text_to_split)
                return
                
            if not separators_list:
                # Force split when no separators left
                step = max(1, self.chunk_size - self.chunk_overlap)
                for i in range(0, len(text_to_split), step):
                    chunks.append(text_to_split[i:i + self.chunk_size])
                return

            sep = separators_list[0]
            parts = text_to_split.split(sep)
            
            current_part = ""
            for part in parts:
                if len(part) > self.chunk_size:
                    if current_part:
                        chunks.append(current_part)
                        current_part = ""
                    recurse(part, separators_list[1:])
                elif len(current_part) + len(part) + (len(sep) if current_part else 0) <= self.chunk_size:
                    current_part += (sep if current_part else "") + part
                else:
                    if current_part:
                        chunks.append(current_part)
                    current_part = part
            if current_part:
                chunks.append(current_part)

        recurse(text, self.separators)
        return [c for c in chunks if c.strip()]

logger = logging.getLogger(__name__)

# Lazy initialization of PaddleOCR
_ocr_instance = None

def get_ocr_instance():
    global _ocr_instance
    if _ocr_instance is None:
        try:
            from paddleocr import PaddleOCR
            _ocr_instance = PaddleOCR(use_angle_cls=True, lang='en', show_log=False)
            logger.info("PaddleOCR successfully initialized.")
        except Exception as e:
            logger.warning(f"PaddleOCR failed to initialize. OCR will be unavailable. Error: {e}")
    return _ocr_instance

class DocumentProcessor:
    """Service for parsing documents, clean text, and chunk content."""
    
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", " ", ""]
        )

    def clean_text(self, text: str) -> str:
        """Removes duplicate whitespaces, control characters, and normalizes spacing."""
        if not text:
            return ""
        text = re.sub(r'[ \t]+', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    def run_ocr_on_page(self, page) -> str:
        """Saves a temporary page image and extracts text using PaddleOCR."""
        ocr = get_ocr_instance()
        if not ocr:
            logger.warning("PaddleOCR instance not available; skipping OCR for page.")
            return ""
            
        temp_img = f"temp_page_{page.number}.png"
        try:
            pix = page.get_pixmap(dpi=150)
            pix.save(temp_img)
            
            result = ocr.ocr(temp_img, cls=True)
            txts = []
            if result and result[0]:
                for line in result[0]:
                    txts.append(line[1][0])
            return " ".join(txts)
        except Exception as e:
            logger.error(f"Error running OCR on PDF page: {e}")
            return ""
        finally:
            if os.path.exists(temp_img):
                try:
                    os.remove(temp_img)
                except Exception:
                    pass

    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extracts text from PDF. Scanned pages are processed via OCR."""
        text_content = []
        try:
            doc = fitz.open(file_path)
            for page in doc:
                page_text = page.get_text()
                if len(page_text.strip()) < 50:
                    ocr_text = self.run_ocr_on_page(page)
                    if ocr_text:
                        page_text = ocr_text
                text_content.append(page_text)
            doc.close()
        except Exception as e:
            logger.error(f"Failed to extract text from PDF: {e}")
            raise e
        return "\n\n".join(text_content)

    def extract_text_from_docx(self, file_path: str) -> str:
        try:
            doc = docx.Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return "\n".join(full_text)
        except Exception as e:
            logger.error(f"Failed to extract text from DOCX: {e}")
            raise e

    def extract_text_from_pptx(self, file_path: str) -> str:
        try:
            prs = pptx.Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception as e:
            logger.error(f"Failed to extract text from PPTX: {e}")
            raise e

    def extract_text_from_image(self, file_path: str) -> str:
        ocr = get_ocr_instance()
        if not ocr:
            raise RuntimeError("PaddleOCR is not available for image text extraction.")
        try:
            result = ocr.ocr(file_path, cls=True)
            txts = []
            if result and result[0]:
                for line in result[0]:
                    txts.append(line[1][0])
            return "\n".join(txts)
        except Exception as e:
            logger.error(f"Failed to extract text from image: {e}")
            raise e

    def extract_text_from_txt(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            logger.error(f"Failed to read TXT file: {e}")
            raise e

    def process_document(self, file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            raw_text = self.extract_text_from_pdf(file_path)
        elif ext == ".docx":
            raw_text = self.extract_text_from_docx(file_path)
        elif ext == ".pptx":
            raw_text = self.extract_text_from_pptx(file_path)
        elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff"]:
            raw_text = self.extract_text_from_image(file_path)
        elif ext == ".txt":
            raw_text = self.extract_text_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
        return self.clean_text(raw_text)

    def chunk_text(self, text: str) -> List[str]:
        return self.text_splitter.split_text(text)
